"""Generate the YouTube presentation deck (.pptx) for the
AI Kubernetes Troubleshooting Agent project.

Run: ../.pptx-venv/bin/python build_deck.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

HERE = Path(__file__).resolve().parent
DOCS = HERE.parent / "docs"

# ── Brand palette (from app/static/style.css) ──────────────────────────────
BG = RGBColor(0x11, 0x13, 0x15)
PANEL = RGBColor(0x1D, 0x23, 0x28)
PANEL2 = RGBColor(0x20, 0x26, 0x2B)
TEXT = RGBColor(0xF4, 0xF7, 0xF8)
MUTED = RGBColor(0x9B, 0xA7, 0xAE)
ACCENT = RGBColor(0x39, 0xC4, 0xA5)
CRIT = RGBColor(0xFF, 0x5A, 0x6B)
HIGH = RGBColor(0xFF, 0x9D, 0x45)
MED = RGBColor(0xF2, 0xD9, 0x5C)
LOW = RGBColor(0x70, 0xC7, 0xFF)
DARKTEXT = RGBColor(0x07, 0x11, 0x0F)

FONT = "Calibri"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    r.fill.solid(); r.fill.fore_color.rgb = BG
    r.line.fill.background()
    r.shadow.inherit = False
    s.shapes._spTree.remove(r._element)
    s.shapes._spTree.insert(2, r._element)
    return s


def accent_bar(s, top=Inches(0.0)):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, Inches(0.18), SH)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False
    return bar


def box(s, l, t, w, h):
    tb = s.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def setrun(p, text, size, color=TEXT, bold=False, font=FONT, italic=False):
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic
    f.name = font; f.color.rgb = color
    return r


def title_slide(kicker, title, subtitle, author):
    s = slide()
    accent_bar(s)
    # kicker
    _, tf = box(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.6))
    p = tf.paragraphs[0]; setrun(p, kicker, 18, ACCENT, bold=True)
    # title
    _, tf = box(s, Inches(0.85), Inches(2.2), Inches(11.8), Inches(2.2))
    p = tf.paragraphs[0]; setrun(p, title, 46, TEXT, bold=True)
    # subtitle
    _, tf = box(s, Inches(0.9), Inches(4.5), Inches(11.3), Inches(1.2))
    p = tf.paragraphs[0]; setrun(p, subtitle, 22, MUTED)
    # author / footer
    _, tf = box(s, Inches(0.9), Inches(6.4), Inches(11.3), Inches(0.6))
    p = tf.paragraphs[0]; setrun(p, author, 16, TEXT, bold=True)
    return s


def header(s, title, kicker=None):
    if kicker:
        _, tf = box(s, Inches(0.6), Inches(0.35), Inches(12), Inches(0.45))
        p = tf.paragraphs[0]; setrun(p, kicker.upper(), 13, ACCENT, bold=True)
        ty = Inches(0.72)
    else:
        ty = Inches(0.45)
    _, tf = box(s, Inches(0.58), ty, Inches(12.2), Inches(0.9))
    p = tf.paragraphs[0]; setrun(p, title, 32, TEXT, bold=True)
    # underline
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.62), Inches(ty.inches + 0.78),
                            Inches(1.6), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = ACCENT
    ln.line.fill.background(); ln.shadow.inherit = False


def bullets(s, items, left=Inches(0.7), top=Inches(1.9), width=Inches(12),
            height=Inches(5), size=20, gap=10):
    tb, tf = box(s, left, top, width, height)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        if isinstance(it, tuple):
            lvl, txt = it
        else:
            lvl, txt = 0, it
        if lvl == 0:
            setrun(p, "▸  ", size, ACCENT, bold=True)
            setrun(p, txt, size, TEXT)
        else:
            p.level = 1
            setrun(p, "•  ", size - 2, MUTED)
            setrun(p, txt, size - 2, MUTED)
    return tb


def card(s, l, t, w, h, fill=PANEL):
    c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.line.color.rgb = PANEL2; c.line.width = Pt(1)
    c.shadow.inherit = False
    try:
        c.adjustments[0] = 0.06
    except Exception:
        pass
    return c


def stat_card(s, l, t, w, h, number, label, ncolor=TEXT):
    card(s, l, t, w, h, PANEL)
    tf = card(s, l, t, w, h).text_frame  # placeholder to keep API simple
    # Instead, add textbox over it
    tb, tf = box(s, l, t + Inches(0.12), w, h)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    setrun(p, number, 30, ncolor, bold=True)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    setrun(p2, label, 12, MUTED, bold=True)


def footer(s, n):
    _, tf = box(s, Inches(11.6), Inches(7.0), Inches(1.6), Inches(0.4))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    setrun(p, f"AI K8s Agent  ·  {n}", 9, MUTED)


# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════
title_slide(
    "AI × DEVOPS · UYGULAMALI EĞİTİM",
    "Kubernetes'te Yapay Zeka Destekli\nSorun Giderme Ajanı",
    "Local bir Kind cluster üzerinde çalışan, kuralları + AI'ı birleştirip\narızaları tespit eden ve onaylı düzeltme öneren bir DevOps asistanı",
    "Hakan Bayraktar  ·  github.com/<repo>  ·  Medium makalesi açıklamada",
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Bu Videoda Ne Var?", "Ajanda")
bullets(s, [
    "Neden AI + DevOps? Bugün operasyonun değiştiği nokta",
    "Problem: Kubernetes sorun gidermek neden zor",
    "Çözüm: Ajan ne yapıyor — mimari ve akış",
    "İki katmanlı zeka: önce deterministik kurallar, sonra AI",
    "Güvenlik & maliyet: maskeleme, rate-limit, insan onayı (human-in-the-loop)",
    "Neden Kind? Gereksinimler ve kod yapısı",
    "CANLI DEMO: arıza üret → tespit et → AI planı → düzelt",
    "Projeyi nasıl genişletirsiniz + repo & Medium makalesi",
], size=21, gap=11)
footer(s, "01")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — AI + DevOps context
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Neden Şimdi? AI + DevOps", "Bağlam")
bullets(s, [
    "Modern sistemler dağıtık, geçici (ephemeral) ve gözlemlenmesi zor",
    "Bir K8s arızasını anlamak için aynı anda bakmak gerekir:",
    (1, "pod durumu · events · loglar · service selector'ları · endpoints · ingress"),
    "Bu bilgi parçalarını birleştirmek zaman ve tecrübe ister",
    "AI burada 'sihir' değil — kanıtları (evidence) hızlıca okuyup yorumlayan bir asistan",
    "Ama AI tek başına yeterli değil: maliyet, halüsinasyon ve güven sorunları var",
    (1, "Bu yüzden önce deterministik kurallar, sonra gerektiğinde AI"),
], size=20, gap=10)
footer(s, "02")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Problem
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Problem: Troubleshooting Dağınık", "Problem")
# left text
bullets(s, [
    "Bir pod CrashLoopBackOff'a girdi. Şimdi ne yapacaksın?",
    (1, "kubectl get pods"),
    (1, "kubectl describe pod ..."),
    (1, "kubectl logs ... --previous"),
    (1, "kubectl get events"),
    (1, "kubectl get endpoints / ingress ..."),
    "Junior bir mühendis için bu zincir göz korkutucu",
    "Senior için bile tekrarlayan, zaman alan bir iş",
], left=Inches(0.7), top=Inches(1.9), width=Inches(7.4), size=19, gap=9)
# right callout card
card(s, Inches(8.5), Inches(2.0), Inches(4.2), Inches(3.6), PANEL)
_, tf = box(s, Inches(8.8), Inches(2.25), Inches(3.7), Inches(3.2))
p = tf.paragraphs[0]; setrun(p, "Ortak arıza tipleri", 16, ACCENT, bold=True)
for t in ["CrashLoopBackOff", "ImagePullBackOff / ErrImagePull",
          "OOMKilled", "CreateContainerConfigError",
          "Pending / FailedScheduling", "Service → No Endpoints",
          "Ingress → Bad Backend"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(6)
    setrun(pp, "› ", 15, ACCENT); setrun(pp, t, 15, TEXT)
footer(s, "03")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Solution overview
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Çözüm: Cluster İçinde Çalışan Bir Ajan", "Çözüm")
bullets(s, [
    "Cluster'ın İÇİNDE bir pod olarak çalışan FastAPI uygulaması",
    "Pod, service, endpoint ve ingress'leri periyodik olarak tarar",
    "Sağlıksız kaynaklar için kanıt toplar: durum, restart, loglar, events",
    "Yerel kural motoru her bulguyu sınıflandırır (rule engine)",
    "Gerekiyorsa maskelenmiş kanıtı Pioneer AI'a gönderir",
    "Bulguları bir web dashboard'da gösterir",
    "Düzeltme HER ZAMAN kullanıcı onayıyla uygulanır (human-in-the-loop)",
], size=20, gap=11)
footer(s, "04")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Core features
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Temel Özellikler", "Özet")
feats = [
    ("Yerel İzleme", "CrashLoop, ImagePull, OOMKilled, ConfigError, Pending tespiti"),
    ("Deterministik Kurallar", "AI çağrısından önce hızlı, ücretsiz analiz"),
    ("Opsiyonel AI", "Karmaşık / bilinmeyen arızalar için Pioneer AI"),
    ("İnsan Onayı", "Hiçbir patch onaysız uygulanmaz"),
    ("Demo Senaryoları", "Kasıtlı bozuk uygulamalarla test"),
    ("Secret Maskeleme", "Log, prompt ve config'de hassas veri gizlenir"),
]
x0, y0 = Inches(0.7), Inches(1.95)
cw, ch = Inches(3.95), Inches(1.65)
gapx, gapy = Inches(0.18), Inches(0.2)
for i, (t, d) in enumerate(feats):
    r, c = divmod(i, 3)
    l = Emu(int(x0) + c * (int(cw) + int(gapx)))
    tp = Emu(int(y0) + r * (int(ch) + int(gapy)))
    card(s, l, tp, cw, ch, PANEL)
    tb, tf = box(s, Emu(int(l) + int(Inches(0.22))), Emu(int(tp) + int(Inches(0.18))),
                 Emu(int(cw) - int(Inches(0.44))), Emu(int(ch) - int(Inches(0.36))))
    p = tf.paragraphs[0]; setrun(p, t, 17, ACCENT, bold=True)
    p2 = tf.add_paragraph(); p2.space_before = Pt(4); setrun(p2, d, 13.5, MUTED)
footer(s, "05")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Architecture (image)
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Mimari", "Nasıl Kurgulandı")
arch = DOCS / "architecture.png"
if arch.exists():
    s.shapes.add_picture(str(arch), Inches(0.7), Inches(1.55),
                         height=Inches(5.4))
_, tf = box(s, Inches(6.7), Inches(2.0), Inches(6.2), Inches(5))
notes = [
    ("Local Machine", "Her şey kendi makinende — bulut maliyeti yok"),
    ("Kind Cluster", "Docker içinde gerçek bir Kubernetes"),
    ("ai-kube-agent ns", "FastAPI ajanı + SQLite cache"),
    ("demo-broken-apps ns", "Kasıtlı bozuk workload'lar"),
    ("Pioneer AI API", "Sadece gerektiğinde, dışarıya HTTPS"),
]
for i, (h, d) in enumerate(notes):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(12)
    setrun(p, h + "  ", 17, ACCENT, bold=True)
    p2 = tf.add_paragraph(); p2.space_after = Pt(2)
    setrun(p2, d, 14, MUTED)
footer(s, "06")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Namespaces & components
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Namespace'ler ve Bileşenler", "Yerleşim")
card(s, Inches(0.7), Inches(1.9), Inches(5.9), Inches(4.7), PANEL)
_, tf = box(s, Inches(0.95), Inches(2.1), Inches(5.4), Inches(4.3))
p = tf.paragraphs[0]; setrun(p, "ai-kube-agent", 18, ACCENT, bold=True, font=MONO)
for t in ["Deployment (FastAPI + dashboard)", "Service (ClusterIP)",
          "ConfigMap (uygulama ayarları)", "Secret (Pioneer API key)",
          "ServiceAccount", "ClusterRole + ClusterRoleBinding (RBAC)"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(7)
    setrun(pp, "› ", 14, ACCENT); setrun(pp, t, 14.5, TEXT)
card(s, Inches(6.85), Inches(1.9), Inches(5.78), Inches(4.7), PANEL)
_, tf = box(s, Inches(7.1), Inches(2.1), Inches(5.3), Inches(4.3))
p = tf.paragraphs[0]; setrun(p, "demo-broken-apps", 18, HIGH, bold=True, font=MONO)
for t in ["crashloop-demo", "imagepull-demo", "oomkilled-demo",
          "bad-config-demo", "service-no-endpoints-demo",
          "ingress-bad-backend-demo", "network-policy-demo", "ai-analysis-demo"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(5)
    setrun(pp, "› ", 14, HIGH); setrun(pp, t, 14, MUTED, font=MONO)
footer(s, "07")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — How it works (scan loop)
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Nasıl Çalışır: Tarama Döngüsü", "Akış")
steps = [
    "Timer ya da UI → Tarama tetiklenir",
    "Scanner → K8s API'den pod, service, ingress, log, event çeker",
    "Rule Engine → deterministik kontroller, bulguları döner",
    "Her bulgu için → fingerprint cache kontrolü (tekrar AI çağrısını önler)",
    "Yeni + AI izinliyse → maskelenmiş kanıt Pioneer AI'a gider",
    "Sonuç + opsiyonel patch önerisi → SQLite'a kaydedilir",
    "Dashboard ve Prometheus metrikleri güncellenir",
]
y = Inches(1.95)
for i, st in enumerate(steps):
    # number chip
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y, Inches(0.5), Inches(0.5))
    chip.fill.solid(); chip.fill.fore_color.rgb = ACCENT
    chip.line.fill.background(); chip.shadow.inherit = False
    ctf = chip.text_frame; ctf.word_wrap = False
    cp = ctf.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    setrun(cp, str(i + 1), 16, DARKTEXT, bold=True)
    _, tf = box(s, Inches(1.4), Emu(int(y) - int(Inches(0.02))), Inches(11.3), Inches(0.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; setrun(p, st, 17.5, TEXT)
    y = Emu(int(y) + int(Inches(0.69)))
footer(s, "08")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Two-layer intelligence
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "İki Katmanlı Zeka", "Tasarım Felsefesi")
# Layer 1
card(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.4), PANEL)
_, tf = box(s, Inches(1.0), Inches(2.25), Inches(5.3), Inches(4))
p = tf.paragraphs[0]; setrun(p, "1 · Rule Engine", 22, ACCENT, bold=True)
for t in ["Deterministik, açıklanabilir", "Hızlı ve ÜCRETSİZ", "Offline çalışır",
          "Bilinen arızaları kapsar", "Her zaman ilk katman"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(9)
    setrun(pp, "✓ ", 16, ACCENT); setrun(pp, t, 16, TEXT)
# Layer 2
card(s, Inches(6.8), Inches(2.0), Inches(5.85), Inches(4.4), PANEL)
_, tf = box(s, Inches(7.1), Inches(2.25), Inches(5.3), Inches(4))
p = tf.paragraphs[0]; setrun(p, "2 · Pioneer AI", 22, LOW, bold=True)
for t in ["Karmaşık / bilinmeyen durumlar", "Kök neden + insan dostu açıklama",
          "Adım adım plan + patch önerisi", "Sadece gerektiğinde tetiklenir",
          "Maliyet ve risk kontrol altında"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(9)
    setrun(pp, "★ ", 16, LOW); setrun(pp, t, 16, TEXT)
footer(s, "09")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Rule engine detail
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Rule Engine: Ne Kontrol Eder?", "Katman 1")
rules = [
    ("CrashLoopBackOff", "restart sayısı, back-off, log ipuçları (connection refused, timeout, panic)"),
    ("ImagePull / ErrImagePull", "image adı, tag, imagePullSecrets"),
    ("OOMKilled", "resource limit'leri, son termination reason"),
    ("Pending / FailedScheduling", "node pressure, taint/toleration, PVC binding"),
    ("CreateContainerConfigError", "eksik ConfigMap / Secret referansı"),
    ("ServiceNoEndpoints", "service selector ↔ pod label eşleşmesi"),
    ("IngressBadBackend", "referans verilen backend service ve port var mı"),
]
y = Inches(1.9)
for name, desc in rules:
    card(s, Inches(0.7), y, Inches(11.95), Inches(0.62), PANEL)
    _, tf = box(s, Inches(0.95), Emu(int(y) + int(Inches(0.04))), Inches(11.5), Inches(0.55))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    setrun(p, f"{name}   ", 15, ACCENT, bold=True, font=MONO)
    setrun(p, desc, 13.5, MUTED)
    y = Emu(int(y) + int(Inches(0.72)))
footer(s, "10")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Pioneer AI integration
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Pioneer AI Entegrasyonu", "Katman 2")
bullets(s, [
    "AI açıkken maskelenmiş kanıt → /v1/chat/completions endpoint'ine gider",
    "Sistem prompt: 'Sen bir Kubernetes SRE'sin. SADECE verilen kanıtı kullan.'",
    "Model çıktısı yapılandırılmış JSON olarak döner:",
], left=Inches(0.7), top=Inches(1.85), width=Inches(12), size=18, gap=9)
card(s, Inches(0.7), Inches(3.3), Inches(11.95), Inches(3.5), PANEL2)
_, tf = box(s, Inches(0.95), Inches(3.45), Inches(11.5), Inches(3.2))
code = [
    '{',
    '  "summary": "...",',
    '  "probable_root_cause": "...",',
    '  "severity": "Critical/High/Medium/Low",',
    '  "recommended_actions": [...],',
    '  "junior_friendly_explanation": "...",',
    '  "action_plan": [...],',
    '  "proposed_fix": { "patch_target": "deployment/v1", "patch_data": {} }',
    '}',
]
for i, ln in enumerate(code):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(2)
    col = ACCENT if ('"' in ln and ':' in ln and 'proposed_fix' not in ln) else TEXT
    setrun(p, ln, 13.5, TEXT, font=MONO)
footer(s, "11")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Security & cost controls
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Güvenlik ve Maliyet Kontrolleri", "Önemli")
items = [
    ("Secret Maskeleme", "Authorization: Bearer…, DATABASE_URL, password=…, token=… dışarı çıkmadan maskelenir"),
    ("Fingerprint Cache", "Aynı bulgu için tekrar tekrar AI çağrısı yapılmaz"),
    ("AI_RATE_LIMIT_PER_SCAN", "Tarama başına AI çağrı sayısı sınırlı (sliding-window)"),
    ("AI_MIN_SEVERITY", "Düşük önem seviyeli bulgular dışarı gönderilmez"),
    ("Cold-start'ta AI kapalı", "Operatör her oturumda AI'ı bilinçli olarak açar"),
]
y = Inches(2.0)
for h, d in items:
    card(s, Inches(0.7), y, Inches(11.95), Inches(0.82), PANEL)
    _, tf = box(s, Inches(0.95), Emu(int(y) + int(Inches(0.08))), Inches(11.5), Inches(0.7))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    setrun(p, h + "   ", 16, ACCENT, bold=True, font=MONO)
    setrun(p, d, 14, MUTED)
    y = Emu(int(y) + int(Inches(0.95)))
footer(s, "12")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Human-in-the-loop remediation
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Düzeltme Akışı: İnsan Döngüde", "Human-in-the-loop")
steps = [
    "Scanner bir bulgu tespit eder",
    "Kullanıcı bulguyu dashboard'da açar",
    "Kullanıcı AI üretimi bir aksiyon planı ister",
    "Üretilen patch incelenir ve ONAYLANIR",
    "Backend patch'i SADECE onaydan sonra uygular",
    "Sonraki tarama iyileşmeyi doğrularsa → 'resolved'",
]
y = Inches(2.0)
for i, st in enumerate(steps):
    chip = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y, Inches(0.55), Inches(0.55))
    chip.fill.solid()
    chip.fill.fore_color.rgb = ACCENT if i != 3 else HIGH
    chip.line.fill.background(); chip.shadow.inherit = False
    cp = chip.text_frame.paragraphs[0]; cp.alignment = PP_ALIGN.CENTER
    setrun(cp, str(i + 1), 17, DARKTEXT, bold=True)
    _, tf = box(s, Inches(1.6), Emu(int(y) - int(Inches(0.02))), Inches(11), Inches(0.6))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    bold = (i == 3)
    setrun(p, st, 18, HIGH if i == 3 else TEXT, bold=bold)
    y = Emu(int(y) + int(Inches(0.74)))
footer(s, "13")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RBAC boundaries
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "RBAC Sınırları", "Güvenli Tasarım")
bullets(s, [
    "ServiceAccount bilinçli olarak kısıtlı yetkilere sahip",
    "Okuma: cluster genelinde pod, service, endpoint, ingress, event",
    "Yazma: sadece seçili demo namespace'indeki seçili kaynaklar",
    "kube-system gibi hassas namespace'ler düzeltme kapsamı DIŞINDA",
    "AI_REMEDIATION_NAMESPACES ile yazma alanı daraltılabilir",
    "AI_REMEDIATION_MODE: sadece-öneri veya uygulanabilir-düzeltme",
], size=20, gap=12)
footer(s, "14")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Why Kind + prerequisites
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Neden Kind? Gereksinimler", "Kurulum")
# left: why kind
card(s, Inches(0.7), Inches(1.95), Inches(5.6), Inches(4.6), PANEL)
_, tf = box(s, Inches(0.95), Inches(2.15), Inches(5.1), Inches(4.2))
p = tf.paragraphs[0]; setrun(p, "Neden Kind?", 18, ACCENT, bold=True)
for t in ["Docker içinde GERÇEK Kubernetes", "Saniyeler içinde kurulur / silinir",
          "Bulut maliyeti ve kimlik yönetimi yok", "Demo için izole, tek seferlik",
          "CI/CD'de de aynı şekilde çalışır", "Öğrenmek ve denemek için ideal"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(9)
    setrun(pp, "✓ ", 15, ACCENT); setrun(pp, t, 15, TEXT)
# right: prereqs
card(s, Inches(6.55), Inches(1.95), Inches(6.1), Inches(4.6), PANEL)
_, tf = box(s, Inches(6.8), Inches(2.15), Inches(5.6), Inches(4.2))
p = tf.paragraphs[0]; setrun(p, "Gereksinimler", 18, ACCENT, bold=True)
for tool, ver in [("Docker / Docker Desktop", "24.0+"), ("Kind", "0.23+"),
                  ("kubectl", "1.29+"), ("Git", "2.0+"), ("Modern tarayıcı", "—"),
                  ("(Opsiyonel) Pioneer API key", "AI için")]:
    pp = tf.add_paragraph(); pp.space_after = Pt(9)
    setrun(pp, "› ", 15, ACCENT)
    setrun(pp, tool + "  ", 15, TEXT, bold=True)
    setrun(pp, ver, 14, MUTED, font=MONO)
footer(s, "15")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Code structure
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Kod Yapısı: Modüller", "Kodlama")
mods = [
    ("main.py", "FastAPI app, route'lar, tarama döngüsü, dashboard"),
    ("scanner.py", "Tarama orkestrasyonu; kural + AI'ı birleştirir"),
    ("rule_engine.py", "Deterministik analizörler (pod / service / ingress)"),
    ("ai_client.py", "Pioneer AI client, prompt'lar, rate limiter"),
    ("k8s_client.py", "Kubernetes API erişimi (list / read / patch)"),
    ("masking.py", "Hassas verilerin maskelenmesi"),
    ("storage.py", "SQLite kalıcılık + ayar persist"),
    ("models.py", "Pydantic modeller, fingerprint üretimi"),
    ("config.py", "Env tabanlı ayarlar (Settings)"),
    ("metrics.py", "Prometheus metrikleri"),
]
y0 = Inches(1.9)
for i, (f, d) in enumerate(mods):
    r, c = divmod(i, 2)
    l = Inches(0.7) if c == 0 else Inches(6.85)
    tp = Emu(int(y0) + r * int(Inches(0.93)))
    card(s, l, tp, Inches(5.95), Inches(0.8), PANEL)
    _, tf = box(s, Emu(int(l) + int(Inches(0.2))), Emu(int(tp) + int(Inches(0.07))),
                Inches(5.6), Inches(0.66))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    setrun(p, f"app/{f}", 14.5, ACCENT, bold=True, font=MONO)
    p2 = tf.add_paragraph(); setrun(p2, d, 12, MUTED)
footer(s, "16")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Dashboard (image)
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Dashboard", "Arayüz")
dash = DOCS / "dashboard.png"
if dash.exists():
    s.shapes.add_picture(str(dash), Inches(0.7), Inches(1.7), width=Inches(8.4))
_, tf = box(s, Inches(9.3), Inches(1.9), Inches(3.6), Inches(5))
for h, d in [("Üst bar", "cluster, kaynak sağlığı, AI durumu"),
             ("Sayaçlar", "Active / Critical / High / Medium / Low"),
             ("Tablo", "aktif & çözülmüş bulgular"),
             ("Detay paneli", "kanıt + AI analizi + plan"),
             ("Run Scan", "manuel tarama tetikle"),
             ("Create Problem", "demo arıza üret")]:
    p = tf.paragraphs[0] if h == "Üst bar" else tf.add_paragraph()
    p.space_after = Pt(11)
    setrun(p, h + "  ", 15, ACCENT, bold=True)
    p2 = tf.add_paragraph(); setrun(p2, d, 13, MUTED)
footer(s, "17")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Live demo plan
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "CANLI DEMO — Akış Planı", "Demo")
bullets(s, [
    "1 · ./scripts/local_test.sh → Kind cluster + ajan + demo'lar tek komutta",
    "2 · Dashboard'ı aç: http://127.0.0.1:18080",
    "3 · 'Create Problem' → Memory Limit Exceeded (OOMKilled) deploy et",
    "4 · 'Run Scan' → bulgunun listeye düşmesini göster",
    "5 · Bulguya tıkla → kanıt + 'Detected by Rules' bölümünü göster",
    "6 · AI'ı aç → 'Generate AI Plan' → kök neden + adım adım plan",
    "7 · Önerilen patch'i incele → onayla → uygula",
    "8 · Tekrar tara → bulgu 'Solved Problems'a geçsin",
], size=18.5, gap=11)
footer(s, "18")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Demo scenarios
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Demo Senaryoları", "Bozuk Uygulamalar")
scen = [
    ("Python Web App Crash Loop", "DB bağlantı hatasıyla sürekli çöker", "HIGH", HIGH),
    ("Go Web App Crash Loop", "2 replica, exit code 1 ile anında çöker", "HIGH", HIGH),
    ("Invalid Container Image", "olmayan tag → ImagePullBackOff", "MEDIUM", MED),
    ("Missing ConfigMap Reference", "eksik ConfigMap → ConfigError", "MEDIUM", MED),
    ("Memory Limit Exceeded", "bellek aşımı → OOMKilled", "HIGH", HIGH),
    ("Service / Ingress hataları", "no-endpoints, bad-backend, network-policy", "MIX", LOW),
]
y = Inches(1.95)
for name, desc, sev, col in scen:
    card(s, Inches(0.7), y, Inches(11.95), Inches(0.75), PANEL)
    _, tf = box(s, Inches(0.95), Emu(int(y) + int(Inches(0.05))), Inches(9.8), Inches(0.66))
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    setrun(p, name + "   ", 16, TEXT, bold=True)
    setrun(p, desc, 13.5, MUTED)
    # severity badge
    bd = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.0), Emu(int(y) + int(Inches(0.19))),
                            Inches(1.3), Inches(0.38))
    bd.fill.solid(); bd.fill.fore_color.rgb = col
    bd.line.fill.background(); bd.shadow.inherit = False
    bp = bd.text_frame.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
    setrun(bp, sev, 11, DARKTEXT, bold=True)
    y = Emu(int(y) + int(Inches(0.83)))
footer(s, "19")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Extending the project
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Projeyi Nasıl Genişletirsiniz?", "Sizin Sıranız")
card(s, Inches(0.7), Inches(2.0), Inches(5.85), Inches(4.3), PANEL)
_, tf = box(s, Inches(1.0), Inches(2.2), Inches(5.3), Inches(3.9))
p = tf.paragraphs[0]; setrun(p, "Yeni AI Sağlayıcı", 18, ACCENT, bold=True)
for t in ["ai_client.py'a yeni client ekle", "config.py'a env değişkenleri ekle",
          "PIONEER_ENDPOINT / PIONEER_MODEL'i yönlendir", "Örn: OpenAI, Bedrock, yerel LLM"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(9)
    setrun(pp, "+ ", 16, ACCENT); setrun(pp, t, 15, TEXT)
card(s, Inches(6.8), Inches(2.0), Inches(5.85), Inches(4.3), PANEL)
_, tf = box(s, Inches(7.1), Inches(2.2), Inches(5.3), Inches(3.9))
p = tf.paragraphs[0]; setrun(p, "Yeni Deterministik Kural", 18, LOW, bold=True)
for t in ["rule_engine.py'daki analizörleri genişlet", "Yeni arıza tipini sınıflandır",
          "tests/test_rule_engine.py'a test ekle", "Örn: HPA, PDB, CertManager kontrolleri"]:
    pp = tf.add_paragraph(); pp.space_after = Pt(9)
    setrun(pp, "+ ", 16, LOW); setrun(pp, t, 15, TEXT)
footer(s, "20")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Repo + Medium
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Kaynaklar: Repo + Makale", "Devamı")
card(s, Inches(0.7), Inches(2.1), Inches(11.95), Inches(1.5), PANEL)
_, tf = box(s, Inches(1.0), Inches(2.35), Inches(11.4), Inches(1.1))
p = tf.paragraphs[0]; setrun(p, "GitHub Repo  ", 18, ACCENT, bold=True)
setrun(p, "— tüm kod, manifest'ler, demo'lar ve scriptler", 15, MUTED)
p2 = tf.add_paragraph(); setrun(p2, "git clone … && ./scripts/local_test.sh", 15, TEXT, font=MONO)
card(s, Inches(0.7), Inches(3.85), Inches(11.95), Inches(1.5), PANEL)
_, tf = box(s, Inches(1.0), Inches(4.1), Inches(11.4), Inches(1.1))
p = tf.paragraphs[0]; setrun(p, "Medium Makalesi  ", 18, ACCENT, bold=True)
setrun(p, "— derinlemesine anlatım, tasarım kararları, ekran görüntüleri", 15, MUTED)
p2 = tf.add_paragraph(); setrun(p2, "Link açıklamada / video sonunda", 15, TEXT)
_, tf = box(s, Inches(0.7), Inches(5.7), Inches(11.95), Inches(1))
p = tf.paragraphs[0]
setrun(p, "İpucu: ", 16, HIGH, bold=True)
setrun(p, "Önce repo'yu çalıştırın, sonra makaleyle derinleşin — kod + yazı birlikte en iyi öğretir.", 16, TEXT)
footer(s, "21")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Takeaways
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s); header(s, "Özet — Akılda Kalsın", "Sonuç")
bullets(s, [
    "AI, DevOps'ta sihir değil; kanıtı hızlı okuyan bir asistan",
    "Önce deterministik kurallar → hız, ücretsizlik, açıklanabilirlik",
    "AI sadece gerektiğinde → maliyet ve risk kontrol altında",
    "Maskeleme + rate-limit + insan onayı = güvenli otomasyon",
    "Kind sayesinde her şey lokalde, maliyetsiz ve tekrar edilebilir",
    "Genişletilebilir mimari: yeni provider, yeni kural eklemek kolay",
], size=21, gap=13)
footer(s, "22")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Outro / CTA
# ════════════════════════════════════════════════════════════════════════
s = slide(); accent_bar(s)
_, tf = box(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.5))
p = tf.paragraphs[0]; setrun(p, "Teşekkürler!", 48, TEXT, bold=True)
_, tf = box(s, Inches(0.95), Inches(3.7), Inches(11.3), Inches(1.5))
p = tf.paragraphs[0]
setrun(p, "Repo'yu klonla, local_test.sh'i çalıştır, kendi kuralını ekle.", 22, MUTED)
p2 = tf.add_paragraph(); p2.space_before = Pt(10)
setrun(p2, "Beğendiysen abone ol · Repo ve Medium linki açıklamada", 18, ACCENT, bold=True)
_, tf = box(s, Inches(0.95), Inches(6.4), Inches(11.3), Inches(0.6))
p = tf.paragraphs[0]; setrun(p, "Hakan Bayraktar — AI Kubernetes Troubleshooting Agent", 15, TEXT, bold=True)

out = HERE / "AI-K8s-Agent-Sunum.pptx"
prs.save(str(out))
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
